"""
Parser for PPTX documents.
"""

import logging
import os
from typing import List
import pptx

from src.parsers.base import BaseParser

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Parser for PPTX documents."""
    
    @staticmethod
    def supported_mime_types() -> List[str]:
        """
        Get the list of MIME types supported by this parser.
        
        Returns:
            List of supported MIME types
        """
        return [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint'
        ]
    
    @classmethod
    def can_parse(cls, mime_type: str) -> bool:
        """
        Check if this parser can handle the given mime type.
        
        Args:
            mime_type: MIME type of the document
            
        Returns:
            True if this parser can handle the mime type, False otherwise
        """
        return mime_type in cls.supported_mime_types() or 'powerpoint' in mime_type.lower()
    
    def parse(self, file_path: str) -> str:
        """
        Parse a PPTX file and extract text content.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Extracted text content
        """
        if not os.path.exists(file_path):
            logger.error(f"PPTX file not found: {file_path}")
            return ""
        
        try:
            text_content = []
            
            # Open the presentation
            presentation = pptx.Presentation(file_path)
            
            # Extract core properties if available
            if hasattr(presentation, 'core_properties'):
                core_props = presentation.core_properties
                if hasattr(core_props, 'title') and core_props.title:
                    text_content.append(f"Title: {core_props.title}")
                if hasattr(core_props, 'author') and core_props.author:
                    text_content.append(f"Author: {core_props.author}")
                text_content.append("")  # Empty line
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                slide_number = i + 1
                text_content.append(f"--- Slide {slide_number} ---")
                
                # Extract slide title
                if slide.shapes.title:
                    title = slide.shapes.title.text
                    text_content.append(f"Title: {title}")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    
                    shape_text = shape.text.strip()
                    if shape_text and shape_text != title:  # Avoid duplicating the title
                        text_content.append(shape_text)
                
                # Add an empty line between slides
                text_content.append("")
            
            return "\n".join(text_content)
            
        except Exception as e:
            logger.exception(f"Error parsing PPTX file {file_path}: {e}")
            return f"Error parsing PPTX file: {str(e)}"
